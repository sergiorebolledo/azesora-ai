import os
import sys
import json
import pandas as pd

# La consola de Windows no usa UTF-8 por defecto y rompe al hacer print() de emojis
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

from langchain_community.document_loaders import PyMuPDFLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from docx import Document as DocumentWord
from pptx import Presentation
from bs4 import BeautifulSoup

EXTENSIONES_GENERICAS = {"docx", "xlsx", "pptx", "json", "html", "htm"}


def inferir_pais(ruta):
    """Infiere el país a partir de la carpeta contenedora (ej. data_source/chile/...)."""
    partes = os.path.normpath(ruta).replace("\\", "/").split("/")
    return "chile" if "chile" in partes else "global"


def extraer_texto_docx(fuente):
    """fuente puede ser una ruta de archivo o un objeto tipo archivo (ej. subida de Streamlit)."""
    doc = DocumentWord(fuente)
    partes = [p.text for p in doc.paragraphs if p.text.strip()]
    for tabla in doc.tables:
        for fila in tabla.rows:
            texto_fila = " | ".join(celda.text.strip() for celda in fila.cells)
            if texto_fila.strip(" |"):
                partes.append(texto_fila)
    return "\n".join(partes)


def extraer_texto_pptx(fuente):
    """Devuelve una lista de (numero_diapositiva, texto) incluyendo notas del orador."""
    prs = Presentation(fuente)
    resultado = []
    for i, slide in enumerate(prs.slides, start=1):
        fragmentos = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                fragmentos.append(shape.text_frame.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            fragmentos.append(f"[Notas del orador]: {slide.notes_slide.notes_text_frame.text}")
        resultado.append((i, "\n".join(fragmentos)))
    return resultado


def extraer_texto_html(html_bruto):
    soup = BeautifulSoup(html_bruto, "html.parser")
    for etiqueta in soup(["script", "style"]):
        etiqueta.decompose()
    return soup.get_text(separator="\n", strip=True)


def extraer_registros_json(json_bruto):
    """Si el JSON es una lista de objetos planos, devuelve una frase por registro.
    En cualquier otro caso, devuelve el JSON completo como texto para dividir en chunks."""
    data = json.loads(json_bruto)
    if isinstance(data, list) and data and all(isinstance(item, dict) for item in data):
        return [". ".join(f"{clave}: {valor}" for clave, valor in item.items()) for item in data]
    return json.dumps(data, ensure_ascii=False, indent=2)

def procesar_archivos_corporativos():
    documentos_finales = []
    
    # Configuramos el divisor de texto (Chunking) profesional
    # Dividirá en bloques de 600 caracteres con 100 de superposición
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=600,
        chunk_overlap=100,
        length_function=len
    )
    
    print("📂 Iniciando extracción y limpieza de documentos...")

    # 1. PROCESAR PDF (Chile)
    pdf_path = "data_source/chile/regulacion_impuestos_chile.pdf"
    if os.path.exists(pdf_path):
        # PyMuPDFLoader carga un Document por página, conservando el número de página en la metadata
        loader = PyMuPDFLoader(pdf_path)
        paginas_pdf = loader.load()

        # Al dividir sobre los Documents (no sobre texto plano), cada chunk hereda la metadata de su página
        chunks_pdf = text_splitter.split_documents(paginas_pdf)
        for chunk in chunks_pdf:
            chunk.metadata.update({
                "pais": "chile",
                "broker": "todos",
                "fuente": "regulacion_impuestos_chile.pdf",
                "pagina": chunk.metadata.get("page", 0) + 1  # PyMuPDFLoader indexa páginas desde 0
            })
            documentos_finales.append(chunk)
        print(f"✅ PDF de Chile procesado. Creados {len(chunks_pdf)} fragmentos con número de página indexado.")

    # 2. PROCESAR MARKDOWN (Global)
    md_path = "data_source/global/interactive_brokers_guide.md"
    if os.path.exists(md_path):
        with open(md_path, "r", encoding="utf-8") as f:
            texto_md = f.read()
        
        chunks_md = text_splitter.split_text(texto_md)
        for chunk in chunks_md:
            doc = Document(
                page_content=chunk,
                metadata={"pais": "global", "broker": "interactive brokers", "fuente": "interactive_brokers_guide.md"}
            )
            documentos_finales.append(doc)
        print(f"✅ Markdown Global procesado. Creados {len(chunks_md)} fragmentos.")

    # 3. PROCESAR CSV (Tabular - Comparativa)
    csv_path = "data_source/comparativa_comisiones.csv"
    if os.path.exists(csv_path):
        df = pd.read_csv(csv_path)
        # Convertimos cada fila del CSV en una frase estructurada comprensible para la IA
        for _, fila in df.iterrows():
            texto_fila = (
                f"El broker {fila['BROKER']} exige un depósito mínimo de {fila['DEPOSITO_MINIMO_USD']} dólares. "
                f"Su comisión por compra es de {fila['COMISION_COMPRA']}. "
                f"El monto mínimo de retiro es de {fila['RETIRO_MINIMO_USD']} dólares y "
                f"¿Tiene soporte en español?: {fila['SOPORTE_ESPANOL']}."
            )
            doc = Document(
                page_content=texto_fila,
                metadata={"pais": "global", "broker": str(fila['BROKER']).lower(), "fuente": "comparativa_comisiones.csv"}
            )
            documentos_finales.append(doc)
        print(f"✅ CSV Comparativo procesado. Creadas {len(df)} filas indexadas.")

    # 4. PROCESAR FORMATOS ADICIONALES (Word, Excel, PowerPoint, JSON, HTML)
    # Escaneamos toda la carpeta data_source/ buscando estas extensiones, sin tocar
    # los archivos ya manejados arriba (PDF/MD/CSV) para no duplicarlos.
    for carpeta_actual, _, archivos in os.walk("data_source"):
        for nombre_archivo in sorted(archivos):
            extension = nombre_archivo.rsplit(".", 1)[-1].lower() if "." in nombre_archivo else ""
            if extension not in EXTENSIONES_GENERICAS:
                continue

            ruta = os.path.join(carpeta_actual, nombre_archivo)
            metadata_base = {"pais": inferir_pais(ruta), "broker": "todos", "fuente": nombre_archivo}
            nuevos_docs = []

            try:
                if extension == "docx":
                    texto = extraer_texto_docx(ruta)
                    for chunk in text_splitter.split_text(texto):
                        nuevos_docs.append(Document(page_content=chunk, metadata=dict(metadata_base)))

                elif extension == "xlsx":
                    hojas = pd.read_excel(ruta, sheet_name=None)
                    for nombre_hoja, df_hoja in hojas.items():
                        for _, fila in df_hoja.iterrows():
                            texto_fila = ". ".join(f"{col}: {val}" for col, val in fila.items())
                            nuevos_docs.append(Document(page_content=texto_fila, metadata=dict(metadata_base, hoja=nombre_hoja)))

                elif extension == "pptx":
                    for num_slide, texto_slide in extraer_texto_pptx(ruta):
                        if not texto_slide.strip():
                            continue
                        for chunk in text_splitter.split_text(texto_slide):
                            nuevos_docs.append(Document(page_content=chunk, metadata=dict(metadata_base, diapositiva=num_slide)))

                elif extension == "json":
                    with open(ruta, "r", encoding="utf-8") as f:
                        contenido = extraer_registros_json(f.read())
                    if isinstance(contenido, list):
                        for texto_registro in contenido:
                            nuevos_docs.append(Document(page_content=texto_registro, metadata=dict(metadata_base)))
                    else:
                        for chunk in text_splitter.split_text(contenido):
                            nuevos_docs.append(Document(page_content=chunk, metadata=dict(metadata_base)))

                elif extension in ("html", "htm"):
                    with open(ruta, "r", encoding="utf-8") as f:
                        texto = extraer_texto_html(f.read())
                    for chunk in text_splitter.split_text(texto):
                        nuevos_docs.append(Document(page_content=chunk, metadata=dict(metadata_base)))

                if nuevos_docs:
                    documentos_finales.extend(nuevos_docs)
                    print(f"✅ {nombre_archivo} procesado ({extension.upper()}). Creados {len(nuevos_docs)} fragmentos.")

            except Exception as e:
                print(f"⚠️ No se pudo procesar {nombre_archivo}: {e}")

    print(f"\n🎉 ¡Procesamiento completado! Total de fragmentos listos para embeddings: {len(documentos_finales)}")
    return documentos_finales

if __name__ == "__main__":
    procesar_archivos_corporativos()